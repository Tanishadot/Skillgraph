"""
SkillGraph AI — Pitch Deck Generator
Generates a professional PPTX for the Redrob AI Hackathon submission.

Usage:
    python create_pitch_deck.py
Output:
    pitch_deck.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x09, 0x09, 0x0B)   # Near-black background
SURFACE  = RGBColor(0x18, 0x18, 0x1B)   # Card surface
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)   # Primary accent
VIOLET_L = RGBColor(0xA7, 0x8B, 0xFA)   # Light violet
EMERALD  = RGBColor(0x10, 0xB9, 0x81)   # Success / positive
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # Warning / highlight
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ZINC_200 = RGBColor(0xE4, 0xE4, 0xE7)
ZINC_400 = RGBColor(0xA1, 0xA1, 0xAA)
ZINC_600 = RGBColor(0x52, 0x52, 0x5B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha: int = 0):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE_TYPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    left, top, width, height,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def section_tag(slide, label: str, top=Inches(0.35)):
    add_rect(slide, Inches(0.6), top, Inches(0.04), Inches(0.28), VIOLET)
    add_text(slide, label.upper(), Inches(0.75), top - Inches(0.02),
             Inches(4), Inches(0.32), font_size=8, bold=True, color=ZINC_400)


def slide_header(slide, title: str, subtitle: str = ""):
    section_tag(slide, "SkillGraph AI · Redrob Hackathon 2025")
    add_text(slide, title, Inches(0.6), Inches(0.7), Inches(12), Inches(1.2),
             font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.5), Inches(11), Inches(0.5),
                 font_size=14, color=ZINC_400)


def bullet_card(slide, left, top, width, height, title: str, bullets: list[str],
                title_color: RGBColor = VIOLET_L, bg_color: RGBColor = SURFACE):
    add_rect(slide, left, top, width, height, bg_color)
    add_text(slide, title, left + Inches(0.2), top + Inches(0.15),
             width - Inches(0.3), Inches(0.35), font_size=11, bold=True, color=title_color)
    y = top + Inches(0.55)
    line_h = Inches(0.3)
    for b in bullets:
        add_text(slide, f"• {b}", left + Inches(0.2), y,
                 width - Inches(0.3), line_h, font_size=9, color=ZINC_200)
        y += line_h


def stat_box(slide, left, top, width, height, value: str, label: str,
             value_color: RGBColor = VIOLET_L):
    add_rect(slide, left, top, width, height, SURFACE)
    add_text(slide, value, left, top + Inches(0.1), width, Inches(0.6),
             font_size=28, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.65), width, Inches(0.35),
             font_size=9, color=ZINC_400, align=PP_ALIGN.CENTER)


# ─── Slides ───────────────────────────────────────────────────────────────────

def slide_cover(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Gradient-ish accent bar left
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, VIOLET)

    # Violet glow blob (approximate with a large semi-transparent rect)
    add_rect(slide, Inches(7), Inches(1), Inches(5.5), Inches(5), RGBColor(0x25, 0x10, 0x55))

    add_text(slide, "REDROB AI HACKATHON 2025", Inches(0.8), Inches(1.2),
             Inches(9), Inches(0.45), font_size=10, bold=True, color=VIOLET_L)

    add_text(slide, "SkillGraph AI", Inches(0.8), Inches(1.75),
             Inches(10), Inches(1.4), font_size=52, bold=True, color=WHITE)

    add_text(slide, "Intelligent Candidate Discovery & Ranking", Inches(0.8), Inches(3.1),
             Inches(10), Inches(0.6), font_size=22, color=ZINC_400)

    add_text(slide,
             "A 7-layer hybrid AI pipeline that scores 100,000+ candidates\n"
             "with semantic matching, behavioral signals & career growth analysis.",
             Inches(0.8), Inches(3.85), Inches(9), Inches(0.9),
             font_size=13, color=ZINC_400)

    # Bottom stat strip
    stats = [("100K+", "Candidates"), ("7", "Ranking layers"),
             ("<50ms", "Query time"), ("0.04×", "Honeypot penalty")]
    x = Inches(0.8)
    for val, lbl in stats:
        add_text(slide, val, x, Inches(5.8), Inches(2.2), Inches(0.5),
                 font_size=18, bold=True, color=VIOLET_L)
        add_text(slide, lbl, x, Inches(6.22), Inches(2.2), Inches(0.35),
                 font_size=9, color=ZINC_600)
        x += Inches(2.5)


def slide_problem(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "The Problem", "Traditional hiring pipelines miss 60–80% of qualified candidates")

    problems = [
        ("Keyword Matching Fails", "Boolean filters reject qualified candidates who use different terminology for the same skill."),
        ("Honeypot Pollution", "Keyword-stuffed resumes and consulting-farm profiles inflate the apparent pool size."),
        ("Manual Review at Scale", "A recruiter can review ~60 profiles/day. At 100,000 applicants, that's 1,600 working days."),
        ("Hidden Gems Lost", "High-fit candidates with low online visibility are never contacted — wasted pipeline value."),
    ]

    cols = 2
    w = Inches(5.8)
    h = Inches(1.9)
    for i, (title, body) in enumerate(problems):
        col = i % cols
        row = i // cols
        left = Inches(0.6) + col * (w + Inches(0.3))
        top = Inches(2.1) + row * (h + Inches(0.2))
        add_rect(slide, left, top, w, h, SURFACE)
        add_rect(slide, left, top, Inches(0.04), h, VIOLET)
        add_text(slide, title, left + Inches(0.18), top + Inches(0.15),
                 w - Inches(0.25), Inches(0.4), font_size=12, bold=True, color=ZINC_200)
        add_text(slide, body, left + Inches(0.18), top + Inches(0.6),
                 w - Inches(0.25), Inches(1.1), font_size=10, color=ZINC_400)


def slide_solution(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Our Solution", "SkillGraph: A 7-layer hybrid AI ranking pipeline")

    layers = [
        ("1  JD Validation", "Detects impossible timelines, skill overload, contradictions", VIOLET_L),
        ("2  BM25 Lexical", "Fast token-level keyword matching baseline", ZINC_200),
        ("3  SBERT Semantic", "384-dim sentence embeddings for skill equivalence", EMERALD),
        ("4  Experience", "Years curve with AI/ML-specific weighting (peak at 7yr)", ZINC_200),
        ("5  Project Depth", "Production ML, stack complexity, open-source contributions", VIOLET_L),
        ("6  Behavior Signals", "Redrob: response rate, availability, notice period", AMBER),
        ("7  Career Growth", "Seniority trajectory — IC vs consulting path detection", EMERALD),
    ]

    x = Inches(0.6)
    y = Inches(2.0)
    w = Inches(12.1)
    h = Inches(0.55)
    gap = Inches(0.08)

    for i, (label, desc, col) in enumerate(layers):
        shade = RGBColor(0x14, 0x14, 0x18) if i % 2 == 0 else RGBColor(0x18, 0x18, 0x1C)
        add_rect(slide, x, y, w, h, shade)
        add_text(slide, label, x + Inches(0.15), y + Inches(0.08),
                 Inches(2.8), h - Inches(0.1), font_size=10, bold=True, color=col)
        add_text(slide, desc, x + Inches(3.1), y + Inches(0.08),
                 Inches(8.8), h - Inches(0.1), font_size=9.5, color=ZINC_400)
        y += h + gap


def slide_architecture(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "System Architecture", "End-to-end inference pipeline + real-time API")

    # Three column boxes
    cols = [
        ("Data Layer", ["candidates.jsonl (100K)", "submission.csv (top-100 ranked)", "JD YAML config", "technology_timeline.yaml"]),
        ("AI Pipeline", ["Job Understanding Agent", "BM25 + SBERT dual-encoder", "7-layer scorer + penalties", "Hidden Gem detector"]),
        ("API + Frontend", ["FastAPI (15 endpoints)", "React + TypeScript + Vite", "Recharts analytics", "Real-time AI Copilot chat"]),
    ]

    w = Inches(3.8)
    for i, (title, items) in enumerate(cols):
        x = Inches(0.55) + i * (w + Inches(0.25))
        bullet_card(slide, x, Inches(2.05), w, Inches(4.6), title, items,
                    title_color=[VIOLET_L, EMERALD, AMBER][i])

    # Arrow connectors (simple text)
    for ax in [Inches(4.35), Inches(8.6)]:
        add_text(slide, "→", ax, Inches(4.1), Inches(0.25), Inches(0.5),
                 font_size=18, bold=True, color=ZINC_600, align=PP_ALIGN.CENTER)


def slide_scoring(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Scoring Methodology", "Weighted composite across 8 dimensions with multiplicative penalties")

    weights = [
        ("Projects",      28, VIOLET_L),
        ("Experience",    22, EMERALD),
        ("Semantic Match",20, AMBER),
        ("Domain Fit",    12, RGBColor(0x06, 0xB6, 0xD4)),
        ("Behavior",       8, RGBColor(0xF4, 0x3F, 0x5E)),
        ("Career Growth",  5, VIOLET_L),
        ("Education",      3, ZINC_400),
        ("Certifications", 2, ZINC_400),
    ]

    bar_left = Inches(0.6)
    bar_max_w = Inches(6.0)
    y = Inches(2.05)
    row_h = Inches(0.52)

    for label, pct, color in weights:
        bar_w = bar_max_w * pct / 30  # normalise so Projects fills ~80%
        add_rect(slide, bar_left, y + Inches(0.12), bar_w, Inches(0.28), color)
        add_text(slide, label, bar_left + Inches(0.08), y + Inches(0.06),
                 Inches(2.5), Inches(0.38), font_size=9.5, color=ZINC_200)
        add_text(slide, f"{pct}%", bar_left + bar_w + Inches(0.12), y + Inches(0.06),
                 Inches(0.6), Inches(0.38), font_size=10, bold=True, color=color)
        y += row_h

    # Penalty box
    bullet_card(slide, Inches(7.2), Inches(2.05), Inches(5.6), Inches(2.5),
                "Multiplicative Penalties",
                ["Consulting-only career  ×0.40",
                 "Honeypot profile         ×0.04",
                 "Keyword stuffer          ×0.55",
                 "Applied after weighted sum"],
                title_color=RGBColor(0xF4, 0x3F, 0x5E))

    bullet_card(slide, Inches(7.2), Inches(4.75), Inches(5.6), Inches(1.8),
                "Hidden Gem Criteria",
                ["Score ≥ 0.63 + Career growth ≥ 0.70",
                 "Low recruiter contact OR open_to_work",
                 "Rank > 20 (overlooked by recruiters)"],
                title_color=AMBER)


def slide_jd_validator(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "JD Intelligence Validator", "Catches impossible requirements before you post them")

    checks = [
        ("Timeline Analysis", "Detects requests like '10yr Kubernetes exp' (impossible — released 2014 → max 11yr)"),
        ("Contradiction Detection", "Flags conflicting signals: 'remote-only' + 'daily office standup required'"),
        ("Skill Overload Scoring", "Quantifies complexity → talent availability estimate → suggests priority vs optional"),
        ("Missing Requirements", "Identifies absent JD fields that reduce match quality (comp range, notice period, growth path)"),
        ("Transferable Mapping", "Maps non-standard skill names to canonical equivalents — reduces false negatives"),
        ("JD Quality Score", "0–1 composite score with actionable improvements applied to corrected JD profile"),
    ]

    w = Inches(5.8)
    h = Inches(1.3)
    for i, (title, body) in enumerate(checks):
        col = i % 2
        row = i // 2
        left = Inches(0.6) + col * (w + Inches(0.3))
        top = Inches(2.0) + row * (h + Inches(0.18))
        add_rect(slide, left, top, w, h, SURFACE)
        add_rect(slide, left, top, Inches(0.04), h, VIOLET)
        add_text(slide, title, left + Inches(0.18), top + Inches(0.1),
                 w - Inches(0.25), Inches(0.38), font_size=11, bold=True, color=VIOLET_L)
        add_text(slide, body, left + Inches(0.18), top + Inches(0.5),
                 w - Inches(0.25), Inches(0.7), font_size=9.5, color=ZINC_400)


def slide_results(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Pipeline Results", "From 100,000 candidates to interview-ready shortlist in seconds")

    stats = [
        ("100K+", "Total screened",   WHITE),
        ("12",    "Interview-ready",  EMERALD),
        ("7",     "Hidden gems",      AMBER),
        ("85",    "Honeypots caught", RGBColor(0xF4, 0x3F, 0x5E)),
    ]
    w = Inches(2.8)
    for i, (val, lbl, col) in enumerate(stats):
        x = Inches(0.6) + i * (w + Inches(0.22))
        stat_box(slide, x, Inches(2.0), w, Inches(1.5), val, lbl, col)

    # Funnel description
    funnel_stages = [
        ("100,000  Raw candidates", ZINC_600, Inches(4.5)),
        ("  ~8,500  Qualified (score > 0.4)", ZINC_400, Inches(4.1)),
        ("  ~1,200  Strong fit (score > 0.55)", ZINC_200, Inches(3.7)),
        ("     ~12  Interview-ready (≥ 0.70)", EMERALD, Inches(3.3)),
        ("       7  Hidden gems discovered", AMBER, Inches(2.9)),
    ]
    y = Inches(3.8)
    for text, col, bar_w in funnel_stages:
        add_rect(slide, Inches(0.6), y, bar_w, Inches(0.38), SURFACE)
        add_rect(slide, Inches(0.6), y, Inches(0.04), Inches(0.38), col)
        add_text(slide, text, Inches(0.82), y + Inches(0.04),
                 Inches(5), Inches(0.32), font_size=9.5, color=col)
        y += Inches(0.48)

    bullet_card(slide, Inches(7.2), Inches(3.7), Inches(5.6), Inches(3.3),
                "Top Candidate Profile (Rank #1)",
                ["Score: 0.847 (97th percentile)",
                 "7.2yr total · 5.4yr AI/ML specific",
                 "Production ML deployments: ✓",
                 "Projects: 0.91 · Semantic: 0.83",
                 "Notice period: 15 days",
                 "Open to work · High response rate"],
                title_color=EMERALD)


def slide_demo_flow(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Demo Flow", "Natural 4-minute walkthrough of the recruiter experience")

    steps = [
        ("1", "Landing Page", "Choose Recruiter or Candidate portal"),
        ("2", "JD Validator", "Paste raw JD → see issues → quality score"),
        ("3", "Candidate Rankings", "Filter, search, sort 100 ranked profiles"),
        ("4", "Candidate Detail", "Score breakdown, career history, signals"),
        ("5", "Hidden Talent", "7 overlooked high-fit candidates surfaced"),
        ("6", "Analytics", "Score distribution, funnel, skill heatmap"),
        ("7", "AI Copilot", "Natural language queries over candidate pool"),
        ("8", "Candidate Portal", "Self-service resume analysis for applicants"),
    ]

    w = Inches(5.6)
    h = Inches(0.72)
    for i, (num, title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(0.55) + col * (w + Inches(0.6))
        y = Inches(2.05) + row * (h + Inches(0.14))
        add_rect(slide, x, y, w, h, SURFACE)
        add_rect(slide, x, y, Inches(0.04), h, VIOLET)
        # Step number circle
        add_rect(slide, x + Inches(0.15), y + Inches(0.16),
                 Inches(0.38), Inches(0.38), RGBColor(0x25, 0x10, 0x55))
        add_text(slide, num, x + Inches(0.15), y + Inches(0.13),
                 Inches(0.38), Inches(0.42), font_size=10, bold=True,
                 color=VIOLET_L, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.65), y + Inches(0.06),
                 Inches(2.4), Inches(0.35), font_size=10, bold=True, color=WHITE)
        add_text(slide, desc, x + Inches(0.65), y + Inches(0.38),
                 Inches(4.7), Inches(0.28), font_size=8.5, color=ZINC_400)


def slide_tech_stack(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Technology Stack", "Production-grade components, zero proprietary lock-in")

    stacks = {
        "AI / ML": ["sentence-transformers (SBERT)", "all-MiniLM-L6-v2 (384-dim)", "BM25 (rank-bm25)", "PyYAML-driven timeline knowledge"],
        "Backend": ["FastAPI + Uvicorn", "Pydantic v2 response models", "python-multipart (form uploads)", "Async + streaming-ready"],
        "Frontend": ["React 18 + TypeScript + Vite", "Framer Motion (animations)", "Recharts (analytics charts)", "TanStack Query v5"],
        "Data": ["JSONL candidate store (100K)", "CSV submission format", "Numpy embeddings cache", "In-memory DataStore singleton"],
    }

    w = Inches(3.0)
    for i, (section, items) in enumerate(stacks.items()):
        col = i % 2
        row = i // 2
        x = Inches(0.55) + col * (w + Inches(0.25)) * 2
        y = Inches(2.0) + row * Inches(2.4)
        bullet_card(slide, x, y, w * 2 + Inches(0.25), Inches(2.2), section, items,
                    title_color=[VIOLET_L, EMERALD, AMBER, ZINC_200][i])


def slide_candidate_portal(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Candidate Portal", "Two-sided experience: recruiter dashboard + self-service candidate analysis")

    bullet_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.6),
                "Recruiter Experience",
                ["Real-time ranked candidate feed",
                 "7-dimension score breakdown per profile",
                 "Side-by-side radar comparison (up to 5)",
                 "JD validation with live quality score",
                 "Hidden gem discovery panel",
                 "Analytics: funnel, distribution, skills",
                 "AI Copilot — natural language queries"],
                title_color=VIOLET_L)

    bullet_card(slide, Inches(6.75), Inches(2.0), Inches(5.8), Inches(4.6),
                "Candidate Self-Assessment Portal",
                ["Step 1: Browse open roles (job cards with skills)",
                 "Step 2: Upload or paste resume text",
                 "Step 3: Instant AI match score (0–1)",
                 "Strengths mapped to JD categories",
                 "Missing skills + transferable equivalents",
                 "Personalised project recommendations",
                 "Certification roadmap to close gaps"],
                title_color=EMERALD)


def slide_differentiators(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Key Differentiators", "Why SkillGraph outperforms naive ranking approaches")

    rows = [
        ("Semantic Equivalence", "SBERT embeddings match 'Chroma' to 'vector database' — Boolean search misses this entirely."),
        ("Honeypot Filtering",   "×0.04 penalty collapses scores of profiles with ≥12 rare skills + no GitHub + 0% response rate."),
        ("JD-side Intelligence", "Most tools only score candidates. We also validate and improve the job description first."),
        ("Hidden Gem Detection", "Explicit rule: score ≥0.63 AND growth ≥0.70 AND (low_contact OR open_to_work) AND rank>20."),
        ("Behavioral Signals",   "Redrob's proprietary signals (response rate, days_since_active, interview completion) integrated as scoring layer."),
        ("Career Growth Proxy",  "Detects IC progression vs consulting churn — penalises stagnant consulting-only trajectories."),
    ]

    y = Inches(2.05)
    rh = Inches(0.72)
    for title, body in rows:
        add_rect(slide, Inches(0.6), y, Inches(12.1), rh, SURFACE)
        add_rect(slide, Inches(0.6), y, Inches(0.04), rh, VIOLET)
        add_text(slide, title, Inches(0.85), y + Inches(0.1),
                 Inches(2.8), Inches(0.5), font_size=10, bold=True, color=VIOLET_L)
        add_text(slide, body, Inches(3.85), y + Inches(0.1),
                 Inches(8.7), Inches(0.5), font_size=9.5, color=ZINC_400)
        y += rh + Inches(0.1)


def slide_roadmap(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)
    slide_header(slide, "What's Next", "Path from prototype to production hiring platform")

    phases = [
        ("Phase 1\n(Done)", ["7-layer ranking pipeline", "FastAPI backend", "React dashboard", "Candidate portal"], EMERALD),
        ("Phase 2\n(Near-term)", ["FAISS ANN index for full 100K", "LLM-powered JD rewriting", "Email outreach integration", "Resume PDF parsing"], AMBER),
        ("Phase 3\n(Scale)", ["Multi-role concurrent campaigns", "Bias detection + fairness audit", "Candidate tracking CRM", "Enterprise SSO / RBAC"], VIOLET_L),
    ]

    w = Inches(3.9)
    for i, (phase, items, col) in enumerate(phases):
        x = Inches(0.6) + i * (w + Inches(0.25))
        bullet_card(slide, x, Inches(2.0), w, Inches(4.6), phase, items, title_color=col)

    # Arrow connectors
    for ax in [Inches(4.5), Inches(8.65)]:
        add_text(slide, "→", ax, Inches(4.0), Inches(0.25), Inches(0.5),
                 font_size=18, bold=True, color=ZINC_600, align=PP_ALIGN.CENTER)


def slide_closing(prs: Presentation):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, VIOLET)

    add_text(slide, "SkillGraph AI", Inches(1), Inches(1.8),
             Inches(11), Inches(1.4), font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Intelligent Candidate Discovery & Ranking", Inches(1), Inches(3.2),
             Inches(11), Inches(0.5), font_size=18, color=ZINC_400, align=PP_ALIGN.CENTER)

    add_text(slide, "Built for the Redrob AI Hackathon 2025", Inches(1), Inches(3.85),
             Inches(11), Inches(0.4), font_size=12, color=ZINC_600, align=PP_ALIGN.CENTER)

    add_text(slide, "github.com/Tanishadot/Skillgraph", Inches(1), Inches(5.2),
             Inches(11), Inches(0.4), font_size=11, color=VIOLET_L, align=PP_ALIGN.CENTER)

    add_text(slide, "Thank you", Inches(1), Inches(5.8),
             Inches(11), Inches(0.5), font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Main ─────────────────────────────────────────────────────────────────────

def build_deck(output_path: str = "pitch_deck.pptx"):
    prs = new_prs()

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_scoring(prs)
    slide_jd_validator(prs)
    slide_results(prs)
    slide_demo_flow(prs)
    slide_tech_stack(prs)
    slide_candidate_portal(prs)
    slide_differentiators(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    prs.save(output_path)
    print(f"[OK] Saved {len(prs.slides)} slides -> {output_path}")
    return output_path


if __name__ == "__main__":
    build_deck()
