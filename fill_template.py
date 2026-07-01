"""
SkillGraph AI — Redrob Idea Submission Template (Filled)
Mirrors the exact template structure from the Redrob x Hack2Skill India.Runs PDF.

Run: python fill_template.py
Out: skillgraph_submission.pptx
"""
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("Run: pip install python-pptx")

# ── Colors from the Redrob template ───────────────────────────────────────────
HDR = RGBColor(0x08, 0x08, 0x12)   # dark header bar
FTR = RGBColor(0x10, 0x22, 0xD0)   # blue footer bar
PUR = RGBColor(0x5B, 0x21, 0xB6)   # purple bullet dot
BLK = RGBColor(0x11, 0x11, 0x11)   # near-black for titles / questions
DGR = RGBColor(0x33, 0x33, 0x33)   # dark gray for answers
MGR = RGBColor(0x55, 0x55, 0x55)   # medium gray for sub-answers
WHT = RGBColor(0xFF, 0xFF, 0xFF)
F   = 'Calibri'

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = WHT

def rct(slide, x, y, w, h, fill=WHT):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def ovl(slide, x, y, w, h, fill=PUR):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def chrome(slide):
    """Apply Redrob template header + white bg + blue footer."""
    bg(slide)
    rct(slide, 0, 0, W, Inches(0.62), fill=HDR)
    # "redrob | H2S" left
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.1), Inches(3.5), Inches(0.44))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run(); r.text = "redrob"; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = WHT; r.font.name = F
    r2 = p.add_run(); r2.text = "  |  H2S"; r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA); r2.font.name = F
    # "INDIA.RUNS" right
    tb2 = slide.shapes.add_textbox(Inches(9.5), Inches(0.08), Inches(3.5), Inches(0.48))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r3 = p2.add_run(); r3.text = "INDIA.RUNS"; r3.font.size = Pt(18)
    r3.font.bold = True; r3.font.italic = True
    r3.font.color.rgb = WHT; r3.font.name = F
    # Blue footer
    rct(slide, 0, H - Inches(0.32), W, Inches(0.32), fill=FTR)

def sec_title(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(9), Inches(0.52))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = BLK; r.font.name = F


def content(slide, items, start_y=Inches(1.38)):
    """
    items: list of tuples:
      ('q', 'Question text')        — bold bullet question
      ('a', 'Answer text')          — indented answer line
      ('s', '')                     — small vertical spacer
    """
    BULLET_X = Inches(0.55)
    Q_X      = Inches(0.78)
    A_X      = Inches(0.98)
    LINE_W   = Inches(12.0)
    A_LINE_W = Inches(11.8)

    y = start_y
    for kind, text in items:
        if kind == 'q':
            # Purple dot
            ovl(slide, BULLET_X, y + Inches(0.08), Inches(0.12), Inches(0.12))
            tb = slide.shapes.add_textbox(Q_X, y, LINE_W, Inches(0.42))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = text
            r.font.size = Pt(12.5); r.font.bold = True
            r.font.color.rgb = BLK; r.font.name = F
            y += Inches(0.44)
        elif kind == 'a':
            tb = slide.shapes.add_textbox(A_X, y, A_LINE_W, Inches(0.36))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = text
            r.font.size = Pt(11); r.font.color.rgb = DGR; r.font.name = F
            y += Inches(0.36)
        elif kind == 'a2':   # sub-answer / continuation, smaller indent
            tb = slide.shapes.add_textbox(A_X + Inches(0.2), y, A_LINE_W - Inches(0.2),
                                          Inches(0.32))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = text
            r.font.size = Pt(10.5); r.font.color.rgb = MGR; r.font.name = F
            y += Inches(0.32)
        elif kind == 's':
            y += Inches(0.15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER (Team Name / Team Leader / Problem Statement)
# ═══════════════════════════════════════════════════════════════════════════════

def s01(prs):
    s = blank(prs)
    bg(s)
    # Approximate the INDIA.RUNS banner at top (dark gradient area)
    rct(s, 0, 0, W, Inches(3.2), fill=HDR)
    # Decorative violet glow (approximate)
    rct(s, 0, 0, Inches(6), Inches(3.2), fill=RGBColor(0x35, 0x08, 0x55))
    rct(s, Inches(7), 0, Inches(6.33), Inches(3.2), fill=RGBColor(0x1A, 0x04, 0x30))
    # redrob | H2S  (top left of banner)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "redrob"; r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = WHT; r1.font.name = F
    r2 = p.add_run(); r2.text = "  |  H2S"; r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB); r2.font.name = F
    # INDIA.RUNS big centered
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(0.85), Inches(11.33), Inches(1.4))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r3 = p2.add_run(); r3.text = "INDIA.RUNS"
    r3.font.size = Pt(64); r3.font.bold = True; r3.font.italic = True
    r3.font.color.rgb = WHT; r3.font.name = F
    # "Build what next India runs on"
    tb3 = s.shapes.add_textbox(Inches(3.0), Inches(2.25), Inches(7.33), Inches(0.55))
    tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    rct(s, Inches(3.2), Inches(2.22), Inches(6.93), Inches(0.55),
        fill=RGBColor(0x08, 0x08, 0x20))
    r4 = p3.add_run(); r4.text = "Build what next India runs on"
    r4.font.size = Pt(14); r4.font.color.rgb = WHT; r4.font.name = F

    # Blue footer on cover
    rct(s, 0, H - Inches(0.32), W, Inches(0.32), fill=FTR)

    # White area below banner: Team info
    fields = [
        ("Team Name :",        "SkillGraph"),
        ("Team Leader Name :", "Tanisha (TanishaDot)"),
        ("Problem Statement :", "Intelligent Candidate Discovery & Ranking Challenge"),
    ]
    for i, (label, value) in enumerate(fields):
        y = Inches(3.4 + i * 0.95)
        # Label
        tb = s.shapes.add_textbox(Inches(0.55), y, Inches(2.8), Inches(0.55))
        tf = tb.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = label
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = BLK; r.font.name = F
        # Value
        tb2 = s.shapes.add_textbox(Inches(3.5), y, Inches(9.3), Inches(0.55))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = value
        r2.font.size = Pt(14); r2.font.color.rgb = DGR; r2.font.name = F


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════

def s02(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Solution Overview")
    content(s, [
        ('q', "What is your proposed solution?"),
        ('a', "SkillGraph AI is a 7-layer hybrid candidate ranking system that combines BM25 sparse "
              "retrieval with SBERT dense semantic embeddings (all-MiniLM-L6-v2, 384-dim) and six "
              "domain-specific signals to rank candidates with sub-second latency."),
        ('a', "A JD Validator improves hiring profiles before ranking, a Hidden Gem detector surfaces "
              "overlooked high-fit talent from lower ranks, and a Candidate Portal enables self-assessment."),
        ('a', "The system is fully explainable — every ranking decision is traceable to structured "
              "resume fields with no LLM dependency in the critical path."),
        ('s', ''),
        ('q', "What differentiates your approach from traditional candidate matching systems?"),
        ('a', "Traditional ATS uses keyword overlap — it rejects candidates who use semantically "
              "equivalent terms. SkillGraph goes beyond keywords in five ways:"),
        ('a2', "1. SBERT semantic embeddings capture true meaning — not surface keyword frequency."),
        ('a2', "2. Career trajectory analysis distinguishes IC growth from consulting-only patterns."),
        ('a2', "3. Behavioral signals from Redrob data (response rate, availability, open_to_work)."),
        ('a2', "4. 3-tier anomaly detection: honeypot x0.04, keyword stuffer x0.55, consulting x0.40."),
        ('a2', "5. Explainable per-candidate reasoning derived from structured fields — no hallucination."),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — JD UNDERSTANDING & CANDIDATE EVALUATION
# ═══════════════════════════════════════════════════════════════════════════════

def s03(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "JD Understanding & Candidate Evaluation")
    content(s, [
        ('q', "What are the key requirements extracted from the JD?"),
        ('a', "Role: Senior AI / ML Engineer  |  Experience: 5+ years total, 3+ years AI/ML"),
        ('a', "Core skills: Python, PyTorch, TensorFlow, MLOps, distributed training"),
        ('a', "Location: India — Bangalore / Hyderabad / Remote  |  Notice period: 60 days ideal"),
        ('a', "Must have: Production ML deployments at scale"),
        ('a', "Nice to have: Open-source contributions, LLM / GenAI experience"),
        ('a', "JD Validator also flags impossible constraints (e.g., 10 years exp in a 5-year-old tool)."),
        ('s', ''),
        ('q', "Which candidate signals are most important? / How does your solution evaluate fit "
              "beyond keyword matching?"),
        ('a', "Beyond keywords, SkillGraph evaluates candidates on these signals:"),
        ('a2', "SBERT semantic similarity — cosine distance to JD embedding, not keyword overlap"),
        ('a2', "Production ML history — real deployed models vs hobby / academic projects"),
        ('a2', "AI/ML-specific years — weighted separately from total career years"),
        ('a2', "Career growth trajectory — IC seniority progression vs consulting-only pattern"),
        ('a2', "Behavioral signals — recruiter response rate, availability, days since last active"),
        ('a2', "Domain alignment — title and role history match to ML engineering career path"),
        ('a2', "Anomaly flags — keyword stuffing detected via BM25 >> SBERT gap"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RANKING METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════

def s04(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Ranking Methodology")
    content(s, [
        ('q', "How does your system retrieve, score, and rank candidates?"),
        ('a2', "1. JD Validator checks quality and builds structured hiring profile"),
        ('a2', "2. BM25 retrieves top candidates from submission pool by keyword relevance"),
        ('a2', "3. SBERT computes dense semantic similarity (resume vs JD embedding)"),
        ('a2', "4. Six domain signals computed: experience, projects, behavior, growth, education, certs"),
        ('a2', "5. Weighted composite score calculated, penalty multipliers applied for anomalies"),
        ('a2', "6. Top-100 candidates returned with score breakdown and natural-language reasoning"),
        ('s', ''),
        ('q', "What models, algorithms, or heuristics are used?"),
        ('a', "SBERT: sentence-transformers all-MiniLM-L6-v2 (384-dim dense embeddings, CPU-friendly)"),
        ('a', "BM25: rank-bm25 library for sparse keyword retrieval"),
        ('a', "Cosine similarity: scikit-learn for SBERT score computation"),
        ('a', "Custom scoring: peak-weighted experience function (maximizes at 7 years, then plateaus)"),
        ('s', ''),
        ('q', "How are multiple candidate signals combined into a final ranking?"),
        ('a', "Weighted sum — Projects 28%  |  Experience 22%  |  Semantic Match 20%  |  "
              "Domain Fit 12%  |  Behavior 8%  |  Career Growth 5%  |  Education 3%  |  Certs 2%"),
        ('a', "Penalty multipliers then applied: Honeypot x0.04  |  Keyword Stuffer x0.55  |  "
              "Consulting Only x0.40"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EXPLAINABILITY & DATA VALIDATION
# ═══════════════════════════════════════════════════════════════════════════════

def s05(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Explainability & Data Validation")
    content(s, [
        ('q', "How are ranking decisions explained?"),
        ('a', "Each candidate has a per-dimension score breakdown with weight and raw value shown to "
              "the recruiter. A natural-language reasoning string is derived deterministically from "
              "the score data — not generated by an LLM."),
        ('a', "Recruiters see: overall score, rank, recommendation label (Strong Match / Good Match / "
              "Borderline), and a per-signal explanation. Score rows are clickable to expand the "
              "dimension explanation, weight, and evidence."),
        ('s', ''),
        ('q', "How do you prevent hallucinations or unsupported justifications?"),
        ('a', "No LLM is used in the critical ranking path. Every score is derived from structured "
              "resume fields using deterministic, reproducible algorithms."),
        ('a', "Every number shown to the recruiter is directly traceable to a specific field in the "
              "candidate data. The AI chat assistant is additive only — it cannot affect rankings."),
        ('s', ''),
        ('q', "How does your solution handle inconsistent, low-quality, or suspicious profiles?"),
        ('a', "3-tier penalty system applied before final ranking:"),
        ('a2', "Honeypot: score >= 0.85 but career_growth < 0.20  ->  final score x0.04"),
        ('a2', "Keyword Stuffer: BM25 score >> SBERT score (gap above threshold)  ->  x0.55"),
        ('a2', "Consulting Only: no IC roles detected in career history  ->  x0.40"),
        ('a', "JD Validator flags impossible constraints in the job description before ranking begins."),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — END-TO-END WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════════

def s06(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "End-to-End Workflow")
    content(s, [
        ('q', "What is the complete workflow from JD input to ranked candidate output?"),
        ('a2', "1.  Recruiter opens SkillGraph dashboard and uploads job description text"),
        ('a2', "2.  JD Validator checks quality, detects impossible constraints, suggests improvements"),
        ('a2', "3.  Validated JD builds a structured hiring profile (skills, experience, location, notice)"),
        ('a2', "4.  BM25 retrieves top candidates from submission pool by keyword relevance"),
        ('a2', "5.  SBERT computes dense semantic similarity between each resume and the JD embedding"),
        ('a2', "6.  Six domain signals computed per candidate (experience, projects, behavior, growth, "
               "education, certifications)"),
        ('a2', "7.  Weighted composite score calculated; penalty multipliers applied for anomalies"),
        ('a2', "8.  Top-100 candidates ranked and returned with per-dimension score breakdown + reasoning"),
        ('a2', "9.  Hidden Gem detector identifies overlooked high-growth candidates from rank > 20"),
        ('a2', "10. Recruiter can compare candidates side-by-side or ask the AI Copilot questions"),
        ('a2', "11. Candidate Portal (parallel): candidates self-assess by uploading resume against JD"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════

def s07(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "System Architecture")
    content(s, [
        ('q', "Three-Layer Decoupled Architecture"),
        ('s', ''),
        ('a', "FRONTEND LAYER — React 18 + TypeScript + Vite + TailwindCSS + Framer Motion + "
              "Recharts + TanStack Query + React Router v6"),
        ('a2', "Pages: Dashboard | Rankings | JD Analysis | Hidden Talent | Analytics | "
               "AI Copilot | Compare | Candidate Portal"),
        ('s', ''),
        ('a', "API LAYER — FastAPI + Python 3.11 + Pydantic v2 + Uvicorn + CORS middleware"),
        ('a2', "Endpoints: /api/candidates | /api/jd/profile | /api/jd/validate | /api/analytics | "
               "/api/candidates/compare | /api/portal/jobs | /api/portal/analyze | /api/chat"),
        ('a2', "CORS enabled | Auto OpenAPI docs | Async routes | FormData + JSON dual support"),
        ('s', ''),
        ('a', "ML ENGINE + DATA LAYER — SBERT all-MiniLM-L6-v2 (384-dim) + rank-bm25 + "
              "scikit-learn + NumPy + DataStore singleton + submission.csv"),
        ('a2', "7-Layer Scoring Engine | Hidden Gem Detector | Honeypot Filter | Embedding cache "
               "(loaded at startup) | In-memory BM25 index | Real-time ranking"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS & PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════

def s08(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Results & Performance")
    content(s, [
        ('q', "What results or insights demonstrate ranking quality?"),
        ('a', "100 candidates ranked with explainable per-candidate reasoning (score breakdown + "
              "natural language rationale)"),
        ('a', "7-layer hybrid pipeline with 8 weighted dimensions captures diverse candidate signals "
              "beyond keyword frequency"),
        ('a', "Hidden Gem detection algorithm implemented: surfaces high-growth candidates from "
              "rank > 20 with low recruiter visibility"),
        ('a', "3-tier anomaly detection filters inflated, keyword-stuffed, and consulting-only "
              "profiles before they reach recruiters"),
        ('a', "JD Validator with quality scoring and improvement suggestions prevents poor-quality "
              "job descriptions from degrading ranking results"),
        ('a', "Per-candidate score breakdown (8 dimensions + penalty) gives recruiters auditable, "
              "justifiable hiring decisions"),
        ('s', ''),
        ('q', "How does your solution meet the challenge's runtime and compute constraints?"),
        ('a', "BM25 index loaded once at startup — O(1) lookup, no per-query re-indexing"),
        ('a', "SBERT embeddings computed and cached at startup — no re-computation per query"),
        ('a', "Cosine similarity is a vectorized NumPy operation — sub-millisecond per candidate"),
        ('a', "Full top-100 ranking is real-time — no queuing, no batch processing, no external calls"),
        ('a', "Model: all-MiniLM-L6-v2 runs efficiently on CPU — no GPU required"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNOLOGIES USED
# ═══════════════════════════════════════════════════════════════════════════════

def s09(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Technologies Used")
    content(s, [
        ('q', "What technologies, frameworks, and tools were used and why were they selected?"),
        ('s', ''),
        ('a', "FRONTEND — React 18, TypeScript, Vite 5, TailwindCSS, Framer Motion, Recharts, "
              "TanStack Query, React Router v6, Lucide Icons"),
        ('a2', "React 18 + hooks for complex stateful UI  |  TypeScript for type-safe API contracts  |  "
               "Vite for fast builds  |  TailwindCSS for utility-first styling with no CSS bloat"),
        ('s', ''),
        ('a', "BACKEND / API — FastAPI, Python 3.11, Pydantic v2, Uvicorn, CORS middleware"),
        ('a2', "FastAPI for native async and auto OpenAPI docs  |  Pydantic v2 for deep schema "
               "validation  |  Async routes for non-blocking I/O  |  FormData + JSON dual support"),
        ('s', ''),
        ('a', "ML / EMBEDDINGS — sentence-transformers (all-MiniLM-L6-v2), rank-bm25, "
              "scikit-learn, NumPy"),
        ('a2', "SBERT chosen for state-of-the-art semantic similarity without GPU dependency  |  "
               "BM25 for efficient sparse retrieval  |  Hybrid BM25+SBERT outperforms either alone"),
        ('s', ''),
        ('a', "DATA — submission.csv candidate pool, DataStore singleton, in-memory BM25 index, "
              "embedding cache"),
        ('a2', "Single-load in-memory design eliminates database overhead and ensures real-time "
               "ranking latency"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SUBMISSION ASSETS
# ═══════════════════════════════════════════════════════════════════════════════

def s10(prs):
    s = blank(prs)
    chrome(s)
    sec_title(s, "Submission Assets")
    content(s, [
        ('q', "GitHub Repository"),
        ('a', "github.com/Tanishadot/Skillgraph"),
        ('a2', "Full source code | README + setup guide | Frontend + Backend + ML pipeline"),
        ('s', ''),
        ('q', "Demo Video"),
        ('a', "4-minute walkthrough: Landing -> Dashboard -> Rankings -> JD Analysis -> "
              "Hidden Talent -> Candidate Portal -> AI Copilot"),
        ('s', ''),
        ('q', "Presentation"),
        ('a', "skillgraph_pitch_deck.pptx  +  skillgraph_pitch_deck.pdf"),
        ('a2', "Generated by python-pptx from real project data (create_submission_deck.py)"),
        ('s', ''),
        ('q', "Submission CSV"),
        ('a', "submission.csv — top-100 ranked candidates with score breakdowns and reasoning"),
        ('s', ''),
        ('q', "Screenshots  +  Architecture Diagram"),
        ('a', "All 8 application screens captured | Three-layer architecture diagram"),
    ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU (template closing slide)
# ═══════════════════════════════════════════════════════════════════════════════

def s11(prs):
    s = blank(prs)
    bg(s)
    # Same dark gradient background as the template closing slide
    rct(s, 0, 0, W, H, fill=HDR)
    rct(s, 0, 0, Inches(5), H, fill=RGBColor(0x35, 0x08, 0x55))
    rct(s, Inches(8), 0, Inches(5.33), H, fill=RGBColor(0x1A, 0x04, 0x30))
    # redrob | H2S
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "redrob"; r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = WHT; r1.font.name = F
    r2 = p.add_run(); r2.text = "  |  H2S"; r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB); r2.font.name = F
    # INDIA.RUNS
    tb2 = s.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11.33), Inches(1.4))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r3 = p2.add_run(); r3.text = "INDIA.RUNS"
    r3.font.size = Pt(64); r3.font.bold = True; r3.font.italic = True
    r3.font.color.rgb = WHT; r3.font.name = F
    # Build what next India runs on
    tb3 = s.shapes.add_textbox(Inches(3), Inches(2.75), Inches(7.33), Inches(0.55))
    tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    rct(s, Inches(3.2), Inches(2.72), Inches(6.93), Inches(0.55),
        fill=RGBColor(0x08, 0x08, 0x20))
    r4 = p3.add_run(); r4.text = "Build what next India runs on"
    r4.font.size = Pt(14); r4.font.color.rgb = WHT; r4.font.name = F
    # THANK YOU
    tb4 = s.shapes.add_textbox(Inches(0), Inches(4.4), W, Inches(1.2))
    tf4 = tb4.text_frame; p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    r5 = p4.add_run(); r5.text = "THANK YOU"
    r5.font.size = Pt(60); r5.font.bold = True
    r5.font.color.rgb = RGBColor(0xC4, 0xB5, 0xFD); r5.font.name = F
    # Blue footer
    rct(s, 0, H - Inches(0.32), W, Inches(0.32), fill=FTR)


# ── BUILD ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    steps = [
        ("01 Cover",             s01),
        ("02 Solution Overview", s02),
        ("03 JD Evaluation",     s03),
        ("04 Ranking",           s04),
        ("05 Explainability",    s05),
        ("06 Workflow",          s06),
        ("07 Architecture",      s07),
        ("08 Results",           s08),
        ("09 Technologies",      s09),
        ("10 Submission Assets", s10),
        ("11 Thank You",         s11),
    ]
    for label, fn in steps:
        try:
            fn(prs)
            print(f"  [OK] {label}")
        except Exception as e:
            print(f"  [ERR] {label}: {e}")
            import traceback; traceback.print_exc()
            raise
    out = Path(__file__).parent / "skillgraph_submission.pptx"
    prs.save(str(out))
    print(f"\nSaved: {out}")
    return out

if __name__ == "__main__":
    build()
